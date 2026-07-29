# ruff: noqa: E501
from __future__ import annotations

import csv
import hashlib
import io

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationDocument
from pptx.slide import Slide
from pptx.util import Inches, Pt

from catora_api.schemas.service_visibility import ServiceVisibilityScorecard


def findings_csv(scorecard: ServiceVisibilityScorecard) -> bytes:
    stream = io.StringIO(newline="")
    writer = csv.writer(stream)
    writer.writerow(
        ["severity", "category", "rule_id", "rule_version", "title", "url", "evidence", "remediation"]
    )
    for item in scorecard.findings:
        writer.writerow(
            [
                item.severity,
                item.category,
                item.rule_id,
                item.rule_version,
                item.title,
                item.url or "",
                item.evidence,
                item.remediation,
            ]
        )
    return stream.getvalue().encode("utf-8")


def content_brief(scorecard: ServiceVisibilityScorecard) -> bytes:
    missing = [item for item in scorecard.questions if item.coverage_state != "supported"]
    lines = [
        "# Service Visibility Content Brief",
        "",
        f"Current score: {scorecard.score_basis_points / 100:.1f}%",
        f"Pages evaluated: {scorecard.page_count}",
        "",
        "## Priority buyer questions",
    ]
    for item in missing[:15]:
        lines.extend(
            [
                f"### {item.question}",
                item.explanation,
                "Suggested action: add a direct, factual answer supported by visible evidence.",
                "",
            ]
        )
    lines.extend(
        [
            "## Guardrails",
            "Use only verified company facts. Do not invent clients, metrics, certifications, prices or locations.",
        ]
    )
    return "\n".join(lines).encode("utf-8")


def _add_title(slide: Slide, title: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.75))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(26)
    paragraph.font.bold = True


def _add_bullets(slide: Slide, lines: list[str]) -> None:
    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(8.45), Inches(5.55))
    frame = body.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(17)
        paragraph.space_after = Pt(8)


def _add_score_slide(
    presentation: PresentationDocument,
    scorecard: ServiceVisibilityScorecard,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, "Catora Service Visibility Audit")
    score_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.55), Inches(3.2), Inches(1.4))
    score = score_box.text_frame.paragraphs[0]
    score.text = f"{scorecard.score_basis_points / 100:.1f}%"
    score.font.size = Pt(46)
    score.font.bold = True
    score.alignment = PP_ALIGN.CENTER
    _add_bullets(
        slide,
        [
            f"Pages evaluated: {scorecard.page_count}",
            f"Findings: {len(scorecard.findings)}",
            f"Buyer questions evaluated: {len(scorecard.questions)}",
            "Evidence-backed diagnostic; no ranking, AI-citation, traffic, lead or revenue guarantee.",
        ],
    )


def executive_pptx(scorecard: ServiceVisibilityScorecard) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)
    _add_score_slide(presentation, scorecard)

    findings_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(findings_slide, "Highest-priority findings")
    finding_lines = [
        f"{item.severity.upper()} — {item.title} ({item.url or 'site-wide'})"
        for item in scorecard.findings[:12]
    ] or ["No deterministic findings were produced for this snapshot."]
    _add_bullets(findings_slide, finding_lines)

    questions_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(questions_slide, "Buyer-question gaps")
    question_lines = [
        f"{item.coverage_state.replace('_', ' ').title()} — {item.question}"
        for item in scorecard.questions
        if item.coverage_state != "supported"
    ][:12] or ["All evaluated buyer questions have supporting public evidence."]
    _add_bullets(questions_slide, question_lines)

    actions_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(actions_slide, "Implementation guardrails")
    _add_bullets(
        actions_slide,
        [
            "Prioritize exact-page fixes linked to the evidence in the CSV.",
            "Use only verified company facts and public proof.",
            "Keep every proposed change human-reviewed; this release does not publish automatically.",
            "Re-scan after approved public changes to produce before-and-after evidence.",
        ],
    )

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def artifact_metadata(content: bytes) -> tuple[str, int]:
    return hashlib.sha256(content).hexdigest(), len(content)
