# ruff: noqa: I001
from __future__ import annotations

import hashlib
import hmac
import io
import time
import uuid

import pptx
import pytest
from pydantic import ValidationError

from catora_api.schemas.service_visibility import (
    SERVICE_VISIBILITY_PROTOCOL_VERSION,
    ServiceVisibilityBridgeBatch,
    ServiceVisibilitySourceCreateRequest,
)
from catora_api.service_visibility.classification import classify_page
from catora_api.service_visibility.crawler import _xml_locations
from catora_api.service_visibility.engine import build_scorecard
from catora_api.service_visibility.extraction import extract_page
from catora_api.service_visibility.questions import (
    QuestionDefinition,
    build_default_questions,
    evaluate_questions,
)
from catora_api.service_visibility.reports import executive_pptx, findings_csv
from catora_api.service_visibility.security import issue_token, verify_signed_body


HTML = """
<html>
  <head>
    <title>Cloud Consulting Services</title>
    <meta name="description" content="Cloud consulting for growing teams">
    <link rel="canonical" href="https://example.com/services/cloud">
    <script type="application/ld+json">
      {"@type":"Organization","name":"Example Co"}
    </script>
  </head>
  <body>
    <h1>Cloud Consulting Services</h1>
    <p>
      Our process starts with discovery and produces measurable project outcomes
      for healthcare teams.
    </p>
    <p>Book a consultation to get started.</p>
    <a href="/case-studies/acme">Case study</a>
  </body>
</html>
"""


def page(html: str = HTML, *, url: str = "https://example.com/services/cloud"):
    return extract_page(url, html)


def test_extracts_and_classifies_service_page() -> None:
    snapshot = page()
    assert snapshot.h1 == "Cloud Consulting Services"
    assert classify_page(snapshot).page_type == "service"
    assert str(snapshot.internal_links[0]) == "https://example.com/case-studies/acme"


def test_source_requires_authorized_https_domain() -> None:
    with pytest.raises(ValidationError, match="HTTPS"):
        ServiceVisibilitySourceCreateRequest.model_validate(
            {
                "name": "Example",
                "startUrl": "http://example.com",
                "connectionMode": "zero_install",
                "authorizedDomainConfirmed": True,
            }
        )


def test_sitemap_parser_rejects_dtd_and_entities() -> None:
    malicious = b'<!DOCTYPE foo [<!ENTITY xxe SYSTEM "file:///etc/passwd">]><urlset></urlset>'
    with pytest.raises(ValueError, match="declarations"):
        _xml_locations(malicious)


def test_default_suite_has_25_deterministic_questions() -> None:
    first = build_default_questions([page()])
    second = build_default_questions([page()])
    assert len(first) == 25
    assert [item.question_hash for item in first] == [item.question_hash for item in second]
    assert {item.question_type for item in first} >= {
        "definition",
        "fit",
        "problem",
        "outcome",
        "process",
        "timing",
        "cost",
        "proof",
        "expertise",
        "security",
        "limitations",
    }


def test_question_evaluation_supports_partial_and_conflicting_states() -> None:
    question = QuestionDefinition(
        position=1,
        question="How is support delivered?",
        question_type="support",
        entity_key=None,
        required_terms=(("support",), ("monitoring",)),
        question_hash="a" * 64,
    )
    partial_page = page(
        HTML.replace(
            "Our process starts with discovery",
            "Our support team starts with discovery",
        )
    )
    partial = evaluate_questions([partial_page], [question])[0]
    assert partial.coverage_state == "partially_supported"

    conflict_page = page(
        HTML.replace(
            "Book a consultation to get started.",
            "Monitoring support is included. Monitoring support is not available after delivery.",
        )
    )
    conflict = evaluate_questions([conflict_page], [question])[0]
    assert conflict.coverage_state == "conflicting"
    assert conflict.conflicting_evidence


def test_scorecard_has_evidence_backed_findings_and_questions() -> None:
    scorecard = build_scorecard(uuid.uuid4(), uuid.uuid4(), [page()])
    assert scorecard.page_count == 1
    assert 0 <= scorecard.score_basis_points <= 10_000
    assert len(scorecard.questions) == 25
    assert all(item.rule_version for item in scorecard.findings)
    assert "answer.missing_process" not in {item.rule_id for item in scorecard.findings}
    assert "evidence.missing_proof" not in {item.rule_id for item in scorecard.findings}


def test_reports_are_valid_editable_pptx_and_csv() -> None:
    scorecard = build_scorecard(uuid.uuid4(), uuid.uuid4(), [page()])
    assert findings_csv(scorecard).startswith(b"severity,category")
    presentation = pptx.Presentation(io.BytesIO(executive_pptx(scorecard)))
    assert len(presentation.slides) == 4
    assert "Catora Service Visibility Audit" in presentation.slides[0].shapes[0].text


def test_bridge_batch_rejects_duplicate_canonicals() -> None:
    snapshot = page()
    with pytest.raises(ValidationError, match="repeat a canonical"):
        ServiceVisibilityBridgeBatch.model_validate(
            {
                "protocolVersion": SERVICE_VISIBILITY_PROTOCOL_VERSION,
                "snapshotId": str(uuid.uuid4()),
                "sequence": 0,
                "complete": True,
                "pages": [
                    snapshot.model_dump(mode="json", by_alias=True),
                    snapshot.model_dump(mode="json", by_alias=True),
                ],
            }
        )


def test_signed_bridge_body_accepts_current_signature_and_rejects_changes() -> None:
    token, digest, _ = issue_token()
    body = b'{"pages":[]}'
    timestamp = str(int(time.time()))
    signature = hmac.new(
        token.encode(),
        timestamp.encode() + b"." + body,
        hashlib.sha256,
    ).hexdigest()
    verify_signed_body(
        token=token,
        expected_token_hash=digest,
        timestamp=timestamp,
        signature=signature,
        body=body,
    )
    with pytest.raises(ValueError, match="signature"):
        verify_signed_body(
            token=token,
            expected_token_hash=digest,
            timestamp=timestamp,
            signature=signature,
            body=body + b" ",
        )


def test_signed_bridge_body_rejects_stale_timestamp() -> None:
    token, digest, _ = issue_token()
    body = b"{}"
    timestamp = "2000000000"
    signature = hmac.new(
        token.encode(),
        timestamp.encode() + b"." + body,
        hashlib.sha256,
    ).hexdigest()
    with pytest.raises(ValueError, match="clock skew"):
        verify_signed_body(
            token=token,
            expected_token_hash=digest,
            timestamp=timestamp,
            signature=signature,
            body=body,
        )
